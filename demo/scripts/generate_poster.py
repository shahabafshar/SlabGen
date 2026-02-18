"""
Generate a 48x36 inch academic poster (.pptx) for the GPSS Research Conference.

Design:
  - ISU cardinal/gold brand, clean white body
  - Three columns: text-heavy | screenshot-heavy | case study + results
  - Multiple Mo2C GUI screenshots showing full workflow
  - Proper image aspect ratios via Pillow

Usage:  python demo/scripts/generate_poster.py
Output: demo/output/GPSS_poster_SlabGen.pptx
"""
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- paths --
script_dir = Path(__file__).parent
project_root = script_dir.parent.parent
output_dir = script_dir.parent / "output"
output_dir.mkdir(parents=True, exist_ok=True)
gui_ss = output_dir / "gui_screenshots"
demo_viz = output_dir

# -- ISU brand colors --
CARDINAL   = RGBColor(0xC8, 0x10, 0x2E)
GOLD       = RGBColor(0xF1, 0xBE, 0x48)
DARK_CARD  = RGBColor(0x7C, 0x0A, 0x1E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xFA, 0xFA, 0xFA)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
SECTION_BG = RGBColor(0xF5, 0xF5, 0xF5)

# -- poster dims --
PW = Inches(48)
PH = Inches(36)
MARGIN = Inches(0.7)


# =====================================================================
# Helpers
# =====================================================================

def img_dims(path, target_width=None, target_height=None):
    """Calculate image dimensions preserving aspect ratio."""
    p = Path(path)
    if not p.exists():
        return (target_width or Inches(4), target_height or Inches(3))
    img = PILImage.open(p)
    ar = img.width / img.height
    if target_width and not target_height:
        w = target_width
        h = int(w / ar)
        return (Emu(w), Emu(h))
    elif target_height and not target_width:
        h = target_height
        w = int(h * ar)
        return (Emu(w), Emu(h))
    elif target_width and target_height:
        box_ar = target_width / target_height
        if ar > box_ar:
            w = target_width
            h = int(w / ar)
        else:
            h = target_height
            w = int(h * ar)
        return (Emu(w), Emu(h))
    return (Emu(img.width * 9525), Emu(img.height * 9525))


def add_image_fitted(slide, path, left, top, max_w, max_h, center_h=True):
    """Add image preserving aspect ratio within a bounding box."""
    p = Path(path)
    if not p.exists():
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_w, max_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = f"[{p.name}]"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = MED_GRAY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return shape

    w, h = img_dims(path, max_w, max_h)
    x_off = (max_w - w) // 2
    y_off = (max_h - h) // 2 if center_h else 0
    return slide.shapes.add_picture(
        str(path), left + x_off, top + y_off, width=w, height=h)


def rect(slide, left, top, w, h, fill, border=None, border_w=Pt(0.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = border_w
    else:
        s.line.fill.background()
    return s


def rrect(slide, left, top, w, h, fill, border=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def tbox(slide, left, top, w, h):
    return slide.shapes.add_textbox(left, top, w, h)


def para(tf, text, size=20, bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
         name="Calibri", after=Pt(6), first=False, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    p.space_after = after
    return p


def heading(tf, text, size=32, color=CARDINAL, first=False):
    return para(tf, text, size=size, bold=True, color=color,
                after=Pt(10), first=first)


def subheading(tf, text, size=24, color=DARK_CARD, first=False):
    return para(tf, text, size=size, bold=True, color=color,
                after=Pt(6), first=first)


def bullet(tf, text, size=22, color=NEAR_BLACK):
    return para(tf, f"\u2022  {text}", size=size, color=color, after=Pt(5))


def caption(tf, text, size=18, color=MED_GRAY, first=False):
    return para(tf, text, size=size, italic=True, color=color,
                align=PP_ALIGN.CENTER, after=Pt(2), first=first)


def section_card(slide, x, y, w, h, accent_color=CARDINAL):
    """Draw a white rounded-rect card with a cardinal left accent bar."""
    rrect(slide, x, y, w, h, WHITE, border=LIGHT_GRAY)
    rect(slide, x + Inches(0.12), y + Inches(0.12),
         Inches(0.14), min(Inches(1.0), h - Inches(0.24)), accent_color)
    return x, y


# =====================================================================
# Build the poster
# =====================================================================

def build_poster():
    prs = Presentation()
    prs.slide_width = PW
    prs.slide_height = PH
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Layout constants
    header_h = Inches(3.2)
    footer_h = Inches(1.0)
    col_gap = Inches(0.45)
    body_top = header_h + Inches(0.5)
    body_w = PW - 2 * MARGIN
    col_w = (body_w - 2 * col_gap) / 3
    body_bot = PH - footer_h - Inches(0.15)
    body_h = body_bot - body_top
    col1_x = MARGIN
    col2_x = col1_x + col_w + col_gap
    col3_x = col2_x + col_w + col_gap

    sec_pad = Inches(0.35)
    text_w = col_w - 2 * sec_pad - Inches(0.15)
    txt_left_offset = sec_pad + Inches(0.15)

    # =================================================================
    # HEADER - cardinal with gold accent line
    # =================================================================
    rect(slide, 0, 0, PW, header_h, CARDINAL)
    rect(slide, 0, header_h, PW, Pt(6), GOLD)

    # Title
    tb = tbox(slide, Inches(1.0), Inches(0.3), Inches(37), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    para(tf, "SlabGen: An Interactive Platform for Automated",
         size=50, bold=True, color=WHITE, after=Pt(2), first=True)
    para(tf, "Surface Slab Generation, Screening, and DFT Workflow Preparation",
         size=50, bold=True, color=WHITE, after=Pt(0))

    # Authors
    tb = tbox(slide, Inches(1.0), Inches(2.1), Inches(37), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    para(tf, "Shahab Afshar  and  Zeinab Hajali Fard",
         size=28, bold=True, color=GOLD, after=Pt(3), first=True)
    para(tf, "Iowa State University of Science and Technology",
         size=24, bold=False, color=RGBColor(0xFF, 0xDD, 0xDD), after=Pt(0))

    # ISU logo area (top right) - white badge with Cyclones logo or text
    logo_path = output_dir / "isu_logo_hires.png"
    if logo_path.exists():
        # White rounded badge behind the logo for contrast on cardinal
        badge_w, badge_h = Inches(5.5), Inches(2.4)
        badge_x, badge_y = Inches(41.5), Inches(0.35)
        rrect(slide, badge_x, badge_y, badge_w, badge_h, WHITE)
        add_image_fitted(slide, logo_path,
                         badge_x + Inches(0.3), badge_y + Inches(0.15),
                         badge_w - Inches(0.6), badge_h - Inches(0.3))
    else:
        tb = tbox(slide, Inches(40.5), Inches(0.6), Inches(6.5), Inches(2.0))
        tf = tb.text_frame
        para(tf, "IOWA STATE", size=40, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, first=True)
        para(tf, "UNIVERSITY", size=32, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)

    # White body background
    rect(slide, 0, header_h + Pt(6), PW,
         PH - header_h - footer_h - Pt(6), OFF_WHITE)

    # =================================================================
    # COLUMN 1 - Motivation / Workflow / Features
    # =================================================================
    y = body_top

    # -- Motivation --
    sec_h = Inches(7.2)
    section_card(slide, col1_x, y, col_w, sec_h)
    tb = tbox(slide, col1_x + txt_left_offset, y + Inches(0.15),
              text_w, sec_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Motivation", first=True)
    bullet(tf, "Surface properties govern catalytic activity, "
               "corrosion resistance, and thin-film growth")
    bullet(tf, "DFT slab calculations require careful setup: "
               "proper terminations, vacuum, and dipole corrections")
    bullet(tf, "Manual slab construction is tedious and error-prone "
               "-- systematic screening of orientations is impractical by hand")
    bullet(tf, "No existing tool integrates structure retrieval, "
               "generation, screening, and DFT prep in one interface")
    para(tf, "", size=6, after=Pt(6))
    para(tf, "SlabGen automates the entire pipeline from "
             "bulk crystal to DFT-ready slab models.",
         size=24, bold=True, color=CARDINAL, after=Pt(0))

    y += sec_h + Inches(0.3)

    # -- Architecture & Workflow --
    sec_h = Inches(10.2)
    section_card(slide, col1_x, y, col_w, sec_h)
    tb = tbox(slide, col1_x + txt_left_offset, y + Inches(0.15),
              text_w, sec_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Architecture & Workflow", first=True)
    para(tf, "Python  +  pymatgen  +  PySide6 (Qt6)  +  matplotlib",
         size=21, bold=True, color=DARK_GRAY, after=Pt(14))

    steps = [
        ("Search", "Query Materials Project by formula or load POSCAR/CIF"),
        ("Configure", "Set Miller indices (h,k,l), vacuum thickness, z-repetitions"),
        ("Generate", "Create slabs with all unique terminations"),
        ("Screen", "Batch-explore symmetrically distinct surfaces"),
        ("Export", "POSCAR files + complete VASP input sets"),
    ]
    for i, (label, desc) in enumerate(steps, 1):
        para(tf, f"{i}.  {label}:  {desc}", size=21, after=Pt(7))

    para(tf, "", size=6, after=Pt(10))
    subheading(tf, "Key Algorithm: oriented_slab_replication")
    para(tf, "Two-stage approach avoids atom-stretching artifacts:",
         size=20, color=DARK_GRAY, after=Pt(4))
    para(tf, "(1) Orient bulk so target (hkl) aligns with z-axis, replicate",
         size=20, color=DARK_GRAY, after=Pt(2))
    para(tf, "(2) Apply SlabGenerator with (001) for vacuum and terminations",
         size=20, color=DARK_GRAY, after=Pt(0))

    y += sec_h + Inches(0.3)

    # -- Key Features (fill remaining) --
    sec_h = body_bot - y
    section_card(slide, col1_x, y, col_w, sec_h)
    tb = tbox(slide, col1_x + txt_left_offset, y + Inches(0.15),
              text_w, sec_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Key Features", first=True)
    features = [
        "Materials Project API integration (search by formula)",
        "Interactive 3D structure viewer (Jmol color scheme)",
        "Batch surface screening with symmetry analysis",
        "Auto dipole correction (LDIPOL, IDIPOL) for slab DFT",
        "Selective dynamics (freeze bottom layers)",
        "CSV export, batch POSCAR export",
        "SLURM job script generation",
    ]
    for f in features:
        bullet(tf, f, size=21)

    # =================================================================
    # COLUMN 2 - GUI Screenshots (visual-heavy)
    # =================================================================
    y = body_top
    img_pad = Inches(0.3)
    img_max_w = col_w - 2 * img_pad

    # -- Section A: Bulk structure + search --
    sec_h = Inches(9.0)
    section_card(slide, col2_x, y, col_w, sec_h)
    tb = tbox(slide, col2_x + txt_left_offset, y + Inches(0.15),
              text_w, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "1. Structure Retrieval", first=True)
    para(tf, "Search Materials Project for Mo\u2082C, select ground-state "
             "mp-1552 (Pbcn, orthorhombic, 12 atoms).",
         size=19, color=MED_GRAY, after=Pt(0))

    img_top = y + Inches(2.2)
    img_h = sec_h - Inches(2.5)
    add_image_fitted(slide, gui_ss / "04_structure_selected.png",
                     col2_x + img_pad, img_top, img_max_w, img_h)

    y += sec_h + Inches(0.3)

    # -- Section B: Slab generation --
    sec_h = Inches(9.0)
    section_card(slide, col2_x, y, col_w, sec_h)
    tb = tbox(slide, col2_x + txt_left_offset, y + Inches(0.15),
              text_w, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "2. Slab Generation", first=True)
    para(tf, "Mo\u2082C (1,1,1) surface: 6 unique terminations, 36 atoms each, "
             "49.20 \u00c5\u00b2 surface area, 15 \u00c5 vacuum.",
         size=19, color=MED_GRAY, after=Pt(0))

    img_top = y + Inches(2.2)
    img_h = sec_h - Inches(2.5)
    add_image_fitted(slide, gui_ss / "07_slab_selected.png",
                     col2_x + img_pad, img_top, img_max_w, img_h)

    y += sec_h + Inches(0.3)

    # -- Section C: Screening results --
    sec_h = Inches(6.8)
    section_card(slide, col2_x, y, col_w, sec_h)
    tb = tbox(slide, col2_x + txt_left_offset, y + Inches(0.15),
              text_w, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "3. Surface Screening", first=True)
    para(tf, "20 terminations across 7 Mo\u2082C surfaces (max Miller index 1). "
             "Color-coded symmetry, sortable table, CSV export.",
         size=19, color=MED_GRAY, after=Pt(0))

    img_top = y + Inches(2.2)
    img_h = sec_h - Inches(2.5)
    add_image_fitted(slide, gui_ss / "09_screening_results.png",
                     col2_x + img_pad, img_top, img_max_w, img_h)

    # Remaining space for DFT screenshot
    y += sec_h + Inches(0.3)
    sec_h = body_bot - y
    section_card(slide, col2_x, y, col_w, sec_h)
    tb = tbox(slide, col2_x + txt_left_offset, y + Inches(0.15),
              text_w, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "4. DFT Input Generation", first=True)
    para(tf, "VASP INCAR with auto dipole correction, KPOINTS "
             "(k_z=1 for slab), POSCAR, SLURM script. Live preview.",
         size=19, color=MED_GRAY, after=Pt(0))

    img_top = y + Inches(2.2)
    img_h = sec_h - Inches(2.5)
    add_image_fitted(slide, gui_ss / "10_dft_dialog.png",
                     col2_x + img_pad, img_top, img_max_w, img_h)

    # =================================================================
    # COLUMN 3 - Case Study + Results + Conclusions
    # =================================================================
    y = body_top

    # -- Mo2C Case Study intro --
    sec_h = Inches(5.5)
    section_card(slide, col3_x, y, col_w, sec_h)
    tb = tbox(slide, col3_x + txt_left_offset, y + Inches(0.15),
              text_w, sec_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Case Study: Mo\u2082C Surfaces", first=True)
    para(tf, "Molybdenum carbide is a promising non-precious-metal "
             "catalyst for hydrogen evolution (HER), CO\u2082 reduction, "
             "and Fischer-Tropsch synthesis. Understanding surface "
             "terminations is critical for rational catalyst design.",
         size=22, after=Pt(12))
    subheading(tf, "Screening Results")
    bullet(tf, "Structure: mp-1552 (Pbcn, orthorhombic, 12 atoms)", size=21)
    bullet(tf, "20 unique terminations across 7 surface orientations", size=21)
    bullet(tf, "8 symmetric + 12 asymmetric terminations identified", size=21)
    bullet(tf, "Surface areas range: 24.63 - 49.20 \u00c5\u00b2", size=21)

    y += sec_h + Inches(0.3)

    # -- Screening visualization grid (from quick_demo) --
    sec_h = Inches(10.5)
    section_card(slide, col3_x, y, col_w, sec_h)
    tb = tbox(slide, col3_x + txt_left_offset, y + Inches(0.15),
              text_w, Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Mo\u2082C Surface Gallery", first=True)
    para(tf, "3D visualization of four representative Mo\u2082C surfaces "
             "generated by SlabGen.",
         size=19, color=MED_GRAY, after=Pt(0))

    img_top = y + Inches(1.8)
    img_h = sec_h - Inches(2.1)
    # Use the 4-panel screening grid from quick_demo
    add_image_fitted(slide, demo_viz / "03_surface_screening.png",
                     col3_x + img_pad, img_top, img_max_w, img_h)

    y += sec_h + Inches(0.3)

    # -- Screening Results Table (real data from CSV) --
    sec_h = Inches(5.5)
    section_card(slide, col3_x, y, col_w, sec_h, accent_color=GOLD)
    tb = tbox(slide, col3_x + txt_left_offset, y + Inches(0.15),
              text_w, sec_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Screening Summary", color=DARK_CARD, first=True)

    # Monospace table header
    para(tf, "Surface   Terms   Area(\u00c5\u00b2)  Symmetric",
         size=22, bold=True, color=DARK_GRAY,
         name="Consolas", after=Pt(6))
    # Real screening data from Mo2C mp-1552 (max_index=1)
    rows = [
        ("(0,0,1)", "2", "24.63", "2 / 2"),
        ("(0,1,0)", "2", "28.62", "1 / 2"),
        ("(1,0,0)", "1", "31.53", "0 / 1"),
        ("(0,1,1)", "4", "37.76", "2 / 4"),
        ("(1,0,1)", "3", "40.01", "2 / 3"),
        ("(1,1,0)", "2", "42.58", "1 / 2"),
        ("(1,1,1)", "6", "49.20", "2 / 6"),
    ]
    for surf, terms, area, sym in rows:
        line = f"{surf:<10}{terms:<8}{area:<10}{sym}"
        para(tf, line, size=20, color=DARK_GRAY,
             name="Consolas", after=Pt(2))

    para(tf, "", size=6, after=Pt(4))
    para(tf, "Total: 20 terminations | 10 symmetric, 10 asymmetric",
         size=19, bold=True, color=CARDINAL, after=Pt(0))

    y += sec_h + Inches(0.3)

    # -- Conclusions & Future Work --
    sec_h = body_bot - y
    section_card(slide, col3_x, y, col_w, sec_h)
    tb = tbox(slide, col3_x + txt_left_offset, y + Inches(0.15),
              text_w, sec_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    heading(tf, "Conclusions & Future Work", first=True)
    bullet(tf, "SlabGen provides an integrated workflow from "
               "bulk crystal to DFT-ready slab models", size=22)
    bullet(tf, "Systematic screening enables rapid surface "
               "exploration with symmetry classification", size=22)
    bullet(tf, "Mo\u2082C case study: 20 terminations across 7 surfaces "
               "characterized in seconds", size=22)
    para(tf, "", size=6, after=Pt(8))
    subheading(tf, "Future Directions")
    bullet(tf, "DFT surface energy calculations and Wulff construction", size=21)
    bullet(tf, "Automated convergence testing workflows", size=21)
    bullet(tf, "Adsorbate placement on generated surfaces", size=21)

    para(tf, "", size=6, after=Pt(10))
    para(tf, "References", size=22, bold=True, color=DARK_GRAY, after=Pt(4))
    para(tf, "[1] Ong et al., Comp. Mater. Sci. 68, 314 (2013) - pymatgen",
         size=16, color=MED_GRAY, after=Pt(2))
    para(tf, "[2] Jain et al., APL Mater. 1, 011002 (2013) - Materials Project",
         size=16, color=MED_GRAY, after=Pt(2))
    para(tf, "[3] Sun & Ceder, Surf. Sci. 617, 53 (2013) - Surface energies",
         size=16, color=MED_GRAY, after=Pt(0))

    # =================================================================
    # FOOTER
    # =================================================================
    rect(slide, 0, PH - footer_h, PW, footer_h, CARDINAL)
    rect(slide, 0, PH - footer_h, PW, Pt(4), GOLD)

    tb = tbox(slide, Inches(1.0), PH - footer_h + Inches(0.15),
              Inches(30), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    para(tf, "GPSS Research Conference  |  Iowa State University  |  "
             "February 20, 2026",
         size=22, color=RGBColor(0xFF, 0xDD, 0xDD), first=True, after=Pt(2))
    para(tf, "Built with Python  |  pymatgen  |  PySide6  |  matplotlib",
         size=18, color=RGBColor(0xEE, 0xAA, 0xAA))

    tb = tbox(slide, Inches(38), PH - footer_h + Inches(0.15),
              Inches(9), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    para(tf, "github.com/shahab-afshar/SlabGen",
         size=20, color=RGBColor(0xFF, 0xDD, 0xDD),
         align=PP_ALIGN.RIGHT, first=True)

    # =================================================================
    # Save
    # =================================================================
    out_path = output_dir / "GPSS_poster_SlabGen.pptx"
    prs.save(str(out_path))
    print(f"Poster saved: {out_path.absolute()}")
    print(f"  Size: 48 x 36 inches (landscape)")
    print(f"\nScreenshots used:")
    print(f"  - 04_structure_selected.png (bulk Mo2C in 3D viewer)")
    print(f"  - 07_slab_selected.png (slab with properties)")
    print(f"  - 09_screening_results.png (screening table)")
    print(f"  - 10_dft_dialog.png (VASP inputs)")
    print(f"  - 03_surface_screening.png (4-panel surface gallery)")
    print(f"\nTo customize in PowerPoint:")
    print(f"  - Replace X.XX surface energies with your DFT results")
    print(f"  - Optionally replace ISU logo with higher-res version")
    print(f"  - Add QR code to footer if desired")


if __name__ == "__main__":
    build_poster()
