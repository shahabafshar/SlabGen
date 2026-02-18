"""
Generate a 12-slide presentation (.pptx) for the GPSS Research Conference oral.

Widescreen 16:9, ISU cardinal/gold branding, embedded Mo2C screenshots.
Can also be printed as a single large sheet for poster use.

Usage:  python demo/scripts/generate_presentation.py
Output: demo/output/GPSS_presentation_SlabGen.pptx
"""
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- paths --
script_dir = Path(__file__).parent
output_dir = script_dir.parent / "output"
output_dir.mkdir(parents=True, exist_ok=True)
gui_ss = output_dir / "gui_screenshots"
demo_viz = output_dir

# -- ISU brand colors --
CARDINAL   = RGBColor(0xC8, 0x10, 0x2E)
GOLD       = RGBColor(0xF1, 0xBE, 0x48)
DARK_CARD  = RGBColor(0x7C, 0x0A, 0x1E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xF8, 0xF8)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
FAINT_CARD = RGBColor(0xFD, 0xF0, 0xF2)

# -- slide dims (widescreen 16:9) --
SW = Inches(13.333)
SH = Inches(7.5)


# =====================================================================
# Helpers
# =====================================================================

def img_dims(path, target_width=None, target_height=None):
    p = Path(path)
    if not p.exists():
        return (target_width or Inches(4), target_height or Inches(3))
    img = PILImage.open(p)
    ar = img.width / img.height
    if target_width and target_height:
        box_ar = target_width / target_height
        if ar > box_ar:
            w = target_width; h = int(w / ar)
        else:
            h = target_height; w = int(h * ar)
        return (Emu(w), Emu(h))
    elif target_width:
        return (Emu(target_width), Emu(int(target_width / ar)))
    elif target_height:
        return (Emu(int(target_height * ar)), Emu(target_height))
    return (Emu(img.width * 9525), Emu(img.height * 9525))


def add_img(slide, path, left, top, max_w, max_h):
    p = Path(path)
    if not p.exists():
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, max_w, max_h)
        s.fill.solid(); s.fill.fore_color.rgb = LIGHT_GRAY
        s.line.fill.background()
        tf = s.text_frame
        tf.paragraphs[0].text = f"[{p.name}]"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = MED_GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return s
    w, h = img_dims(path, max_w, max_h)
    x_off = (max_w - w) // 2
    y_off = (max_h - h) // 2
    return slide.shapes.add_picture(str(path), left + x_off, top + y_off,
                                    width=w, height=h)


def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def tbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def para(tf, text, size=20, bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
         name="Calibri", after=Pt(6), first=False, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = color; p.font.name = name
    p.alignment = align; p.space_after = after
    return p


def bullet(tf, text, size=22, color=NEAR_BLACK):
    return para(tf, f"\u2022  {text}", size=size, color=color, after=Pt(8))


def numbered(tf, num, text, size=22, color=NEAR_BLACK):
    return para(tf, f"{num}.   {text}", size=size, color=color, after=Pt(8))


# =====================================================================
# Slide templates
# =====================================================================

def slide_header(slide, title, subtitle=None):
    """Cardinal bar at top with slide title."""
    rect(slide, 0, 0, SW, Inches(1.1), CARDINAL)
    rect(slide, 0, Inches(1.1), SW, Pt(4), GOLD)
    tb = tbox(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.85))
    tf = tb.text_frame
    para(tf, title, size=32, bold=True, color=WHITE, first=True, after=Pt(0))
    if subtitle:
        tb2 = tbox(slide, Inches(0.6), Inches(0.7), Inches(11), Inches(0.35))
        tf2 = tb2.text_frame
        para(tf2, subtitle, size=18, color=RGBColor(0xFF, 0xDD, 0xDD),
             first=True, after=Pt(0))


def slide_number(slide, num, total=12):
    """Slide number in bottom right."""
    tb = tbox(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = tb.text_frame
    para(tf, f"{num}/{total}", size=12, color=MED_GRAY,
         align=PP_ALIGN.RIGHT, first=True, after=Pt(0))


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# =====================================================================
# Build presentation
# =====================================================================

def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ------------------------------------------------------------------
    # SLIDE 1: Title
    # ------------------------------------------------------------------
    s = new_slide(prs)
    rect(s, 0, 0, SW, SH, CARDINAL)
    rect(s, 0, Inches(5.8), SW, Pt(5), GOLD)

    # ISU logo badge
    logo = output_dir / "isu_logo_hires.png"
    if logo.exists():
        badge_w, badge_h = Inches(2.8), Inches(1.3)
        rrect(s, Inches(10.0), Inches(0.4), badge_w, badge_h, WHITE)
        add_img(s, logo, Inches(10.2), Inches(0.45), badge_w - Inches(0.4), badge_h - Inches(0.1))

    tb = tbox(s, Inches(0.8), Inches(1.0), Inches(11.5), Inches(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "SlabGen", size=56, bold=True, color=GOLD, first=True, after=Pt(8))
    para(tf, "An Interactive Platform for Automated Surface Slab",
         size=34, bold=True, color=WHITE, after=Pt(2))
    para(tf, "Generation, Screening, and DFT Workflow Preparation",
         size=34, bold=True, color=WHITE, after=Pt(0))

    tb = tbox(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Shahab Afshar  and  Zeinab Hajali Fard",
         size=24, bold=True, color=GOLD, first=True, after=Pt(4))
    para(tf, "Iowa State University of Science and Technology",
         size=20, color=RGBColor(0xFF, 0xDD, 0xDD), after=Pt(0))

    tb = tbox(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.6))
    tf = tb.text_frame
    para(tf, "GPSS Research Conference  |  February 20, 2026",
         size=18, color=RGBColor(0xEE, 0xAA, 0xAA), first=True, after=Pt(0))

    # ------------------------------------------------------------------
    # SLIDE 2: Motivation
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Motivation")
    slide_number(s, 2)

    tb = tbox(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Why surface slab modeling matters",
         size=28, bold=True, color=CARDINAL, first=True, after=Pt(20))

    bullet(tf, "Surface properties govern catalytic activity, corrosion "
               "resistance, and thin-film growth", size=24)
    bullet(tf, "DFT slab calculations require careful construction: "
               "proper terminations, vacuum thickness, and dipole corrections", size=24)
    bullet(tf, "Manual slab setup is tedious and error-prone -- "
               "systematic screening of multiple orientations is impractical by hand", size=24)
    bullet(tf, "No existing tool integrates structure retrieval, slab generation, "
               "screening, and DFT input preparation in a single GUI", size=24)

    para(tf, "", size=10, after=Pt(20))
    # highlight box
    rrect(s, Inches(1.5), Inches(5.8), Inches(10.0), Inches(0.9), FAINT_CARD, border=CARDINAL)
    tb2 = tbox(s, Inches(1.8), Inches(5.9), Inches(9.4), Inches(0.7))
    tf2 = tb2.text_frame
    para(tf2, "SlabGen automates the entire pipeline from bulk crystal "
              "to DFT-ready slab models.", size=24, bold=True, color=CARDINAL,
         align=PP_ALIGN.CENTER, first=True, after=Pt(0))

    # ------------------------------------------------------------------
    # SLIDE 3: Background -- Surface Science Concepts
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Background", "Key concepts in computational surface science")
    slide_number(s, 3)

    # Left column
    tb = tbox(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Surface Slab Model", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(12))
    bullet(tf, "Miller indices (h,k,l) define the crystallographic "
               "plane for cutting", size=20)
    bullet(tf, "Terminations: different atomic layers exposed at the surface "
               "(same orientation, different cuts)", size=20)
    bullet(tf, "Vacuum layer separates periodic slab images in DFT", size=20)
    bullet(tf, "Symmetric slabs avoid spurious dipole moments; "
               "asymmetric slabs need dipole correction", size=20)

    para(tf, "", size=8, after=Pt(12))
    para(tf, "DFT for Surfaces", size=26, bold=True, color=CARDINAL, after=Pt(12))
    bullet(tf, "Surface energy: \u03b3 = (E_slab - n\u00d7E_bulk) / (2A)", size=20)
    bullet(tf, "Wulff construction: equilibrium crystal shape from \u03b3 values", size=20)
    bullet(tf, "VASP inputs: INCAR (ISIF=2), KPOINTS (k_z=1), POSCAR", size=20)

    # Right column -- bulk structure image
    rrect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), OFF_WHITE, border=LIGHT_GRAY)
    add_img(s, demo_viz / "01_bulk_structure.png",
            Inches(7.0), Inches(1.7), Inches(5.4), Inches(5.0))

    # ------------------------------------------------------------------
    # SLIDE 4: SlabGen Architecture
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "SlabGen Architecture")
    slide_number(s, 4)

    tb = tbox(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Technology Stack", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(12))
    bullet(tf, "Python + pymatgen (materials analysis)", size=22)
    bullet(tf, "PySide6 / Qt6 (cross-platform GUI)", size=22)
    bullet(tf, "matplotlib (embedded 3D visualization)", size=22)
    bullet(tf, "Materials Project API (structure database)", size=22)

    para(tf, "", size=8, after=Pt(16))
    para(tf, "Modular Design", size=26, bold=True, color=CARDINAL, after=Pt(12))

    modules = [
        ("core/slab_generator.py", "Slab generation + oriented replication"),
        ("core/screening.py", "Batch Miller index screening"),
        ("core/dft_inputs.py", "VASP input file generation"),
        ("core/visualization.py", "3D structure plotting"),
        ("ui/main_window.py", "Primary GUI with integrated viewer"),
    ]
    for mod, desc in modules:
        para(tf, f"{mod}", size=18, bold=True, color=DARK_GRAY,
             name="Consolas", after=Pt(1))
        para(tf, f"    {desc}", size=18, color=MED_GRAY, after=Pt(6))

    # Right side -- app screenshot
    rrect(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5), OFF_WHITE, border=LIGHT_GRAY)
    add_img(s, gui_ss / "01_initial_state.png",
            Inches(7.2), Inches(1.7), Inches(5.4), Inches(5.1))

    # ------------------------------------------------------------------
    # SLIDE 5: Key Algorithm
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Key Algorithm: oriented_slab_replication")
    slide_number(s, 5)

    tb = tbox(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Two-stage approach avoids atom-stretching artifacts",
         size=26, bold=True, color=CARDINAL, first=True, after=Pt(20))

    # Stage boxes
    rrect(s, Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.0), FAINT_CARD, border=CARDINAL)
    tb2 = tbox(s, Inches(1.1), Inches(2.7), Inches(5.0), Inches(1.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    para(tf2, "Stage 1: Orient", size=24, bold=True, color=CARDINAL,
         first=True, after=Pt(8))
    para(tf2, "Rotate bulk crystal so target (h,k,l) plane aligns with "
              "the z-axis, then replicate along z to desired thickness.",
         size=20, color=DARK_GRAY, after=Pt(0))

    rrect(s, Inches(7.0), Inches(2.6), Inches(5.5), Inches(2.0), FAINT_CARD, border=CARDINAL)
    tb3 = tbox(s, Inches(7.3), Inches(2.7), Inches(5.0), Inches(1.8))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    para(tf3, "Stage 2: Generate", size=24, bold=True, color=CARDINAL,
         first=True, after=Pt(8))
    para(tf3, "Apply SlabGenerator with (0,0,1) on oriented structure "
              "to add vacuum and enumerate all unique terminations.",
         size=20, color=DARK_GRAY, after=Pt(0))

    # Arrow between boxes
    tb4 = tbox(s, Inches(6.2), Inches(3.2), Inches(0.9), Inches(0.6))
    tf4 = tb4.text_frame
    para(tf4, "\u27a1", size=36, color=GOLD, align=PP_ALIGN.CENTER,
         first=True, after=Pt(0))

    # Why this matters
    tb5 = tbox(s, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.0))
    tf5 = tb5.text_frame; tf5.word_wrap = True
    para(tf5, "Why two stages?", size=24, bold=True, color=DARK_GRAY,
         first=True, after=Pt(10))
    bullet(tf5, "Direct SlabGenerator on arbitrary (h,k,l) can produce "
                "distorted atomic positions for complex unit cells", size=20)
    bullet(tf5, "Pre-orienting + z-replication ensures physical atom spacing "
                "is preserved before vacuum is applied", size=20)
    bullet(tf5, "Enables correct slab models for any crystal system "
                "(cubic, orthorhombic, hexagonal, ...)", size=20)

    # ------------------------------------------------------------------
    # SLIDE 6: Demo -- Structure Retrieval
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Demo: Structure Retrieval from Materials Project")
    slide_number(s, 6)

    # Text (left)
    tb = tbox(s, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Search & Select", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(14))
    numbered(tf, 1, "Enter chemical formula (Mo\u2082C)", size=21)
    numbered(tf, 2, "Query Materials Project API", size=21)
    numbered(tf, 3, "Browse results: ID, space group, crystal system, atoms", size=21)
    numbered(tf, 4, "Select structure to load into 3D viewer", size=21)

    para(tf, "", size=8, after=Pt(14))
    para(tf, "Selected: mp-1552", size=24, bold=True, color=CARDINAL, after=Pt(8))
    bullet(tf, "Alpha-Mo\u2082C (Pbcn, orthorhombic)", size=20)
    bullet(tf, "12 atoms per unit cell", size=20)
    bullet(tf, "Ground-state polymorph (E_hull = 0)", size=20)

    # Screenshot (right)
    add_img(s, gui_ss / "04_structure_selected.png",
            Inches(5.6), Inches(1.4), Inches(7.2), Inches(5.5))

    # ------------------------------------------------------------------
    # SLIDE 7: Demo -- Slab Generation
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Demo: Slab Generation -- Mo\u2082C (1,1,1)")
    slide_number(s, 7)

    # Screenshot (left, bigger)
    add_img(s, gui_ss / "06_slabs_generated.png",
            Inches(0.4), Inches(1.4), Inches(7.8), Inches(5.7))

    # Text (right)
    tb = tbox(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Parameters", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(12))
    bullet(tf, "Miller indices: (1, 1, 1)", size=21)
    bullet(tf, "Z repetitions: 3", size=21)
    bullet(tf, "Vacuum: 15 \u00c5", size=21)
    bullet(tf, "All terminations: enabled", size=21)

    para(tf, "", size=8, after=Pt(14))
    para(tf, "Results", size=26, bold=True, color=CARDINAL, after=Pt(12))
    bullet(tf, "6 unique terminations", size=21)
    bullet(tf, "36 atoms per slab", size=21)
    bullet(tf, "49.20 \u00c5\u00b2 surface area", size=21)
    bullet(tf, "2 symmetric + 4 asymmetric", size=21)

    # ------------------------------------------------------------------
    # SLIDE 8: Demo -- Slab Inspection
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Demo: Slab Inspection & 3D Visualization")
    slide_number(s, 8)

    # Screenshot
    add_img(s, gui_ss / "07_slab_selected.png",
            Inches(0.4), Inches(1.4), Inches(7.8), Inches(5.7))

    # Properties callout (right)
    tb = tbox(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Slab Properties", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(14))

    props = [
        ("Formula:", "Mo\u2082C"),
        ("Atoms:", "36"),
        ("Surface area:", "49.20 \u00c5\u00b2"),
        ("Slab thickness:", "16.59 \u00c5"),
        ("Symmetric:", "Yes (shift 0.0)"),
    ]
    for label, val in props:
        para(tf, f"{label}  {val}", size=21, after=Pt(6))

    para(tf, "", size=8, after=Pt(14))
    para(tf, "3D Viewer Features", size=26, bold=True, color=CARDINAL, after=Pt(12))
    bullet(tf, "Interactive rotation and zoom", size=20)
    bullet(tf, "Jmol color scheme (element-based)", size=20)
    bullet(tf, "Lattice box with axis labels", size=20)
    bullet(tf, "Updates on slab selection", size=20)

    # ------------------------------------------------------------------
    # SLIDE 9: Demo -- Surface Screening
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Demo: Systematic Surface Screening")
    slide_number(s, 9)

    # Screenshot
    add_img(s, gui_ss / "09_screening_results.png",
            Inches(0.4), Inches(1.4), Inches(7.2), Inches(5.7))

    # Text (right)
    tb = tbox(s, Inches(8.0), Inches(1.5), Inches(4.8), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Batch Screening", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(14))
    bullet(tf, "All symmetrically distinct Miller indices up to max index", size=20)
    bullet(tf, "Runs in background thread (GUI stays responsive)", size=20)
    bullet(tf, "Progress bar with live updates", size=20)

    para(tf, "", size=8, after=Pt(14))
    para(tf, "Mo\u2082C Results (max index 1)", size=24, bold=True,
         color=CARDINAL, after=Pt(10))
    bullet(tf, "7 surface orientations", size=20)
    bullet(tf, "20 unique terminations", size=20)
    bullet(tf, "Color-coded: green = symmetric, yellow = asymmetric", size=20)

    para(tf, "", size=8, after=Pt(14))
    para(tf, "Export Options", size=24, bold=True, color=CARDINAL, after=Pt(10))
    bullet(tf, "CSV table export", size=20)
    bullet(tf, "Batch POSCAR export (all slabs)", size=20)
    bullet(tf, "Load any slab into main viewer", size=20)

    # ------------------------------------------------------------------
    # SLIDE 10: Screening Summary Table
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Screening Results: Mo\u2082C Surface Characterization")
    slide_number(s, 10)

    # Table - use text boxes with monospace for alignment
    # Header row
    table_x = Inches(0.8)
    table_y = Inches(1.8)
    col_widths = [Inches(2.2), Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.5)]
    row_h = Inches(0.5)

    headers = ["Surface", "Terminations", "Area (\u00c5\u00b2)", "Symmetric", "Asymmetric"]
    data_rows = [
        ("(0,0,1)", "2", "24.63", "2", "0"),
        ("(0,1,0)", "2", "28.62", "1", "1"),
        ("(1,0,0)", "1", "31.53", "0", "1"),
        ("(0,1,1)", "4", "37.76", "2", "2"),
        ("(1,0,1)", "3", "40.01", "2", "1"),
        ("(1,1,0)", "2", "42.58", "1", "1"),
        ("(1,1,1)", "6", "49.20", "2", "4"),
    ]

    # Header background
    rect(s, table_x, table_y, sum(w for w in col_widths), row_h, CARDINAL)
    x = table_x
    for i, hdr in enumerate(headers):
        tb = tbox(s, x, table_y, col_widths[i], row_h)
        tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, hdr, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, first=True, after=Pt(0))
        x += col_widths[i]

    # Data rows
    for r, row in enumerate(data_rows):
        ry = table_y + (r + 1) * row_h
        bg = OFF_WHITE if r % 2 == 0 else WHITE
        rect(s, table_x, ry, sum(w for w in col_widths), row_h, bg, border=LIGHT_GRAY)
        x = table_x
        for c, val in enumerate(row):
            tb = tbox(s, x, ry, col_widths[c], row_h)
            tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            clr = DARK_GRAY
            if c == 3 and int(val) > 0:
                clr = RGBColor(0x1B, 0x7A, 0x2B)  # green
            elif c == 4 and int(val) > 0:
                clr = RGBColor(0xB8, 0x86, 0x0B)  # amber
            para(tf, val, size=20, color=clr, bold=(c == 0),
                 align=PP_ALIGN.CENTER, first=True, after=Pt(0))
            x += col_widths[c]

    # Totals row
    ry = table_y + 8 * row_h
    rect(s, table_x, ry, sum(w for w in col_widths), row_h, FAINT_CARD, border=CARDINAL)
    totals = ["Total", "20", "--", "10", "10"]
    x = table_x
    for c, val in enumerate(totals):
        tb = tbox(s, x, ry, col_widths[c], row_h)
        tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, val, size=20, bold=True, color=CARDINAL,
             align=PP_ALIGN.CENTER, first=True, after=Pt(0))
        x += col_widths[c]

    # Note below table
    tb = tbox(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6))
    tf = tb.text_frame
    para(tf, "Structure: alpha-Mo\u2082C (mp-1552, Pbcn, orthorhombic, 12 atoms)   |   "
             "Max Miller index: 1   |   Z reps: 3   |   Vacuum: 15 \u00c5",
         size=16, color=MED_GRAY, align=PP_ALIGN.CENTER, first=True, after=Pt(0))

    # ------------------------------------------------------------------
    # SLIDE 11: Mo2C Surface Gallery
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Mo\u2082C Surface Gallery")
    slide_number(s, 11)

    # Full-width 4-panel screening grid
    add_img(s, demo_viz / "03_surface_screening.png",
            Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.7))

    # ------------------------------------------------------------------
    # SLIDE 12: Demo -- DFT Inputs
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Demo: DFT Input Generation")
    slide_number(s, 12)

    # Screenshot
    add_img(s, gui_ss / "10_dft_dialog.png",
            Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.7))

    # Text (right)
    tb = tbox(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Generated VASP Inputs", size=26, bold=True, color=CARDINAL,
         first=True, after=Pt(14))
    bullet(tf, "INCAR with slab-optimized settings (ISIF=2, ISMEAR=0)", size=21)
    bullet(tf, "Auto dipole correction (LDIPOL, IDIPOL=3, DIPOL from COM)", size=21)
    bullet(tf, "KPOINTS with k_z = 1 for slab geometry", size=21)
    bullet(tf, "POSCAR + POTCAR.spec", size=21)
    bullet(tf, "SLURM job submission script", size=21)

    para(tf, "", size=8, after=Pt(14))
    para(tf, "Advanced Options", size=24, bold=True, color=CARDINAL, after=Pt(10))
    bullet(tf, "Selective dynamics: freeze bottom N layers", size=20)
    bullet(tf, "Adjustable ENCUT, k-density, EDIFFG", size=20)
    bullet(tf, "Live INCAR and KPOINTS preview tabs", size=20)

    # ------------------------------------------------------------------
    # SLIDE 13 (bonus): Conclusions & Future Work
    # ------------------------------------------------------------------
    s = new_slide(prs)
    slide_header(s, "Conclusions & Future Work")
    slide_number(s, 13, total=13)

    # Two-column layout
    tb = tbox(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Conclusions", size=28, bold=True, color=CARDINAL,
         first=True, after=Pt(14))
    bullet(tf, "SlabGen provides an integrated workflow from "
               "bulk crystal to DFT-ready slab models", size=22)
    bullet(tf, "Systematic screening enables rapid surface exploration "
               "with symmetry classification", size=22)
    bullet(tf, "Mo\u2082C case study: 20 terminations across 7 surfaces "
               "characterized in seconds", size=22)
    bullet(tf, "Open-source, cross-platform (Python + PySide6)", size=22)

    tb = tbox(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Future Directions", size=28, bold=True, color=CARDINAL,
         first=True, after=Pt(14))
    bullet(tf, "DFT surface energy calculations and Wulff construction", size=22)
    bullet(tf, "Automated convergence testing workflows", size=22)
    bullet(tf, "Adsorbate placement on generated surfaces", size=22)
    bullet(tf, "Support for heterostructure interface modeling", size=22)

    para(tf, "", size=10, after=Pt(20))
    para(tf, "References", size=22, bold=True, color=DARK_GRAY, after=Pt(8))
    para(tf, "[1] Ong et al., Comp. Mater. Sci. 68, 314 (2013)",
         size=16, color=MED_GRAY, after=Pt(2))
    para(tf, "[2] Jain et al., APL Mater. 1, 011002 (2013)",
         size=16, color=MED_GRAY, after=Pt(2))
    para(tf, "[3] Sun & Ceder, Surf. Sci. 617, 53 (2013)",
         size=16, color=MED_GRAY, after=Pt(0))

    # Thank you bar at bottom
    rrect(s, Inches(1.5), Inches(6.3), Inches(10.0), Inches(0.8), CARDINAL)
    tb = tbox(s, Inches(1.5), Inches(6.35), Inches(10.0), Inches(0.7))
    tf = tb.text_frame
    para(tf, "Thank you!   |   Questions?",
         size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
         first=True, after=Pt(0))

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    out_path = output_dir / "GPSS_presentation_SlabGen.pptx"
    prs.save(str(out_path))
    print(f"Presentation saved: {out_path.absolute()}")
    print(f"  {len(prs.slides)} slides, 16:9 widescreen")
    print(f"\nSlide overview:")
    titles = [
        "Title", "Motivation", "Background", "Architecture",
        "Key Algorithm", "Demo: Structure Retrieval", "Demo: Slab Generation",
        "Demo: Slab Inspection", "Demo: Surface Screening",
        "Screening Results Table", "Mo2C Surface Gallery",
        "Demo: DFT Inputs", "Conclusions & Future Work",
    ]
    for i, t in enumerate(titles, 1):
        print(f"  {i:2d}. {t}")


if __name__ == "__main__":
    build()
